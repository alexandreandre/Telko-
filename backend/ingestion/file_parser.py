"""
Parseur de fichiers multi-format.
Extrait le texte brut depuis des fichiers PDF, DOCX, XLSX, PPTX et TXT.
Retourne des Documents LangChain avec metadata complètes (source, filename,
page, file_type, last_modified). Fichier illisible ou vide → warning + [].
"""

import logging
import os
from datetime import datetime, timezone
from pathlib import Path

from langchain_core.documents import Document

logger = logging.getLogger(__name__)

# Extensions supportées
_SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx",
    ".jpg", ".jpeg", ".png",
    ".txt", ".md",
    ".csv", ".xlsx",
    ".json", ".html",
}


# ---------------------------------------------------------------------------
# Helpers metadata
# ---------------------------------------------------------------------------

def _base_metadata(file_path: str) -> dict:
    """
    Construit le socle de metadata commun à tous les parseurs.
    last_modified est converti en ISO-8601 UTC.
    """
    p = Path(file_path)
    try:
        mtime = os.path.getmtime(file_path)
        last_modified = datetime.fromtimestamp(mtime, tz=timezone.utc).isoformat()
    except OSError:
        last_modified = ""
    return {
        "source": str(p),
        "filename": p.name,
        "file_type": p.suffix.lower().lstrip("."),
        "last_modified": last_modified,
    }


def _make_doc(content: str, page: int | str, base_meta: dict) -> Document:
    """Crée un Document LangChain avec le contenu et les metadata fusionnés."""
    return Document(
        page_content=content,
        metadata={**base_meta, "page": page},
    )


# ---------------------------------------------------------------------------
# Parseurs par format
# ---------------------------------------------------------------------------

def _parse_pdf(file_path: str, base_meta: dict) -> list[Document]:
    """Extrait le texte page par page avec PyMuPDF (fitz)."""
    try:
        import fitz  # pymupdf
    except ImportError as exc:
        raise ImportError("pymupdf est requis pour parser les PDF : pip install pymupdf") from exc

    docs: list[Document] = []
    try:
        with fitz.open(file_path) as pdf:
            for page_num, page in enumerate(pdf, start=1):
                text = page.get_text("text").strip()
                if text:
                    docs.append(_make_doc(text, page_num, base_meta))
    except Exception as exc:
        logger.warning("PDF illisible '%s' : %s", file_path, exc)
        return []

    if not docs:
        logger.warning("PDF vide ou sans texte extractible : '%s'.", file_path)
    return docs


def _parse_docx(file_path: str, base_meta: dict) -> list[Document]:
    """
    Extrait les paragraphes d'un fichier Word (.docx).
    Les paragraphes consécutifs sont regroupés par blocs de 20 pour limiter
    le nombre de Documents créés tout en conservant la granularité.
    """
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        raise ImportError("python-docx est requis : pip install python-docx") from exc

    try:
        docx = DocxDocument(file_path)
    except Exception as exc:
        logger.warning("DOCX illisible '%s' : %s", file_path, exc)
        return []

    paragraphs = [p.text.strip() for p in docx.paragraphs if p.text.strip()]
    if not paragraphs:
        logger.warning("DOCX vide ou sans texte : '%s'.", file_path)
        return []

    # Regroupement par blocs de 20 paragraphes → un Document par bloc
    block_size = 20
    docs: list[Document] = []
    for block_idx, start in enumerate(range(0, len(paragraphs), block_size), start=1):
        block_text = "\n".join(paragraphs[start : start + block_size])
        docs.append(_make_doc(block_text, block_idx, base_meta))
    return docs


def _parse_pptx(file_path: str, base_meta: dict) -> list[Document]:
    """
    Extrait le texte de chaque slide d'une présentation PowerPoint (.pptx).
    Un Document par slide, la page correspond au numéro de slide.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError("python-pptx est requis : pip install python-pptx") from exc

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        logger.warning("PPTX illisible '%s' : %s", file_path, exc)
        return []

    docs: list[Document] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            docs.append(_make_doc("\n".join(texts), slide_num, base_meta))

    if not docs:
        logger.warning("PPTX vide ou sans texte : '%s'.", file_path)
    return docs


def _parse_image(file_path: str, base_meta: dict) -> list[Document]:
    """
    Applique l'OCR Tesseract (lang='fra+eng') sur une image JPG/PNG.
    Retourne un unique Document avec le texte extrait (page=1).
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise ImportError(
            "pytesseract et Pillow sont requis : pip install pytesseract Pillow"
        ) from exc

    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang="fra+eng").strip()
    except Exception as exc:
        logger.warning("Image illisible '%s' : %s", file_path, exc)
        return []

    if not text:
        logger.warning("OCR n'a extrait aucun texte de '%s'.", file_path)
        return []

    return [_make_doc(text, 1, base_meta)]


def _parse_txt(file_path: str, base_meta: dict) -> list[Document]:
    """
    Lit un fichier texte brut en UTF-8 (fallback latin-1).
    Retourne un unique Document (page=1).
    """
    try:
        try:
            text = Path(file_path).read_text(encoding="utf-8").strip()
        except UnicodeDecodeError:
            text = Path(file_path).read_text(encoding="latin-1").strip()
    except Exception as exc:
        logger.warning("TXT illisible '%s' : %s", file_path, exc)
        return []

    if not text:
        logger.warning("Fichier TXT vide : '%s'.", file_path)
        return []

    return [_make_doc(text, 1, base_meta)]


def _parse_md(file_path: str, base_meta: dict) -> list[Document]:
    """Lit un fichier Markdown en UTF-8 (fallback latin-1). Retourne un unique Document."""
    try:
        try:
            text = Path(file_path).read_text(encoding="utf-8").strip()
        except UnicodeDecodeError:
            text = Path(file_path).read_text(encoding="latin-1").strip()
    except Exception as exc:
        logger.warning("MD illisible '%s' : %s", file_path, exc)
        return []

    if not text:
        logger.warning("Fichier MD vide : '%s'.", file_path)
        return []

    return [_make_doc(text, 1, base_meta)]


def _parse_csv(file_path: str, base_meta: dict) -> list[Document]:
    """
    Parse un CSV avec la stdlib.
    Chaque ligne devient "col1: val1 | col2: val2". Tout le fichier = un Document.
    """
    import csv

    try:
        try:
            raw = Path(file_path).read_text(encoding="utf-8")
        except UnicodeDecodeError:
            raw = Path(file_path).read_text(encoding="latin-1")
    except Exception as exc:
        logger.warning("CSV illisible '%s' : %s", file_path, exc)
        return []

    lines: list[str] = []
    reader = csv.DictReader(raw.splitlines())
    for row in reader:
        entry = " | ".join(f"{k}: {v}" for k, v in row.items() if v is not None)
        if entry.strip():
            lines.append(entry)

    if not lines:
        logger.warning("CSV vide ou sans données : '%s'.", file_path)
        return []

    return [_make_doc("\n".join(lines), 1, base_meta)]


def _parse_xlsx(file_path: str, base_meta: dict) -> list[Document]:
    """
    Parse un fichier Excel avec openpyxl.
    Chaque feuille = un Document séparé avec metadata sheet_name.
    """
    try:
        import openpyxl
    except ImportError as exc:
        raise ImportError("openpyxl est requis pour parser les XLSX : pip install openpyxl") from exc

    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    except Exception as exc:
        logger.warning("XLSX illisible '%s' : %s", file_path, exc)
        return []

    docs: list[Document] = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(h) if h is not None else "" for h in rows[0]]
        lines: list[str] = []
        for row in rows[1:]:
            entry = " | ".join(
                f"{headers[i]}: {v}"
                for i, v in enumerate(row)
                if v is not None and str(v).strip()
            )
            if entry.strip():
                lines.append(entry)
        if lines:
            sheet_meta = {**base_meta, "sheet_name": sheet.title}
            docs.append(Document(page_content="\n".join(lines), metadata={**sheet_meta, "page": sheet.title}))

    wb.close()
    if not docs:
        logger.warning("XLSX vide ou sans données : '%s'.", file_path)
    return docs


def _parse_json(file_path: str, base_meta: dict) -> list[Document]:
    """
    Parse un fichier JSON.
    Si le JSON est une liste → un Document par élément.
    Sinon → un unique Document avec le JSON indenté.
    """
    import json

    try:
        text = Path(file_path).read_text(encoding="utf-8")
        data = json.loads(text)
    except Exception as exc:
        logger.warning("JSON illisible '%s' : %s", file_path, exc)
        return []

    if isinstance(data, list):
        docs: list[Document] = []
        for i, item in enumerate(data, start=1):
            content = json.dumps(item, ensure_ascii=False, indent=2)
            if content.strip():
                docs.append(_make_doc(content, i, base_meta))
        if not docs:
            logger.warning("JSON liste vide : '%s'.", file_path)
        return docs

    content = json.dumps(data, ensure_ascii=False, indent=2)
    if not content.strip():
        logger.warning("JSON vide : '%s'.", file_path)
        return []
    return [_make_doc(content, 1, base_meta)]


def _parse_html(file_path: str, base_meta: dict) -> list[Document]:
    """
    Parse un fichier HTML avec html.parser de la stdlib.
    Supprime toutes les balises et décode les entités HTML.
    Retourne le texte brut en un unique Document.
    """
    from html.parser import HTMLParser
    import html as html_module

    class _TextExtractor(HTMLParser):
        _SKIP_TAGS = {"script", "style", "head"}

        def __init__(self) -> None:
            super().__init__()
            self._parts: list[str] = []
            self._skip = 0

        def handle_starttag(self, tag: str, attrs: list) -> None:
            if tag in self._SKIP_TAGS:
                self._skip += 1

        def handle_endtag(self, tag: str) -> None:
            if tag in self._SKIP_TAGS and self._skip:
                self._skip -= 1

        def handle_data(self, data: str) -> None:
            if not self._skip:
                stripped = data.strip()
                if stripped:
                    self._parts.append(stripped)

        def get_text(self) -> str:
            return "\n".join(self._parts)

    try:
        try:
            raw = Path(file_path).read_text(encoding="utf-8")
        except UnicodeDecodeError:
            raw = Path(file_path).read_text(encoding="latin-1")
    except Exception as exc:
        logger.warning("HTML illisible '%s' : %s", file_path, exc)
        return []

    extractor = _TextExtractor()
    extractor.feed(raw)
    text = html_module.unescape(extractor.get_text()).strip()

    if not text:
        logger.warning("HTML vide ou sans texte : '%s'.", file_path)
        return []

    return [_make_doc(text, 1, base_meta)]


# ---------------------------------------------------------------------------
# API publique
# ---------------------------------------------------------------------------

_PARSERS = {
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".jpg":  _parse_image,
    ".jpeg": _parse_image,
    ".png":  _parse_image,
    ".txt":  _parse_txt,
    ".md":   _parse_md,
    ".csv":  _parse_csv,
    ".xlsx": _parse_xlsx,
    ".json": _parse_json,
    ".html": _parse_html,
}


def parse_file(file_path: str) -> list[Document]:
    """
    Parse un fichier et retourne une liste de Documents LangChain.

    Formats supportés : .pdf, .docx, .pptx, .jpg, .jpeg, .png, .txt, .md, .csv, .xlsx, .json, .html

    Chaque Document contient :
      - page_content : texte extrait
      - metadata     : {source, filename, page, file_type, last_modified}

    Args:
        file_path: Chemin absolu ou relatif vers le fichier.

    Returns:
        Liste de Documents. Retourne [] si :
          - le format n'est pas supporté
          - le fichier est illisible ou introuvable
          - aucun texte n'a pu être extrait
    """
    p = Path(file_path)

    if not p.exists():
        logger.warning("Fichier introuvable : '%s'.", file_path)
        return []

    ext = p.suffix.lower()
    if ext not in _PARSERS:
        logger.warning(
            "Format non supporté '%s' pour '%s'. Extensions acceptées : %s.",
            ext,
            file_path,
            ", ".join(sorted(_SUPPORTED_EXTENSIONS)),
        )
        return []

    base_meta = _base_metadata(file_path)
    parser = _PARSERS[ext]

    logger.info("Parsing '%s' (type=%s)…", p.name, base_meta["file_type"])
    docs = parser(file_path, base_meta)
    logger.info("'%s' → %d Document(s) extraits.", p.name, len(docs))
    return docs


def batch_parse(folder_path: str) -> list[Document]:
    """
    Parse tous les fichiers supportés d'un dossier (non récursif).

    Les fichiers dont le format n'est pas supporté sont ignorés silencieusement.
    Les fichiers illisibles génèrent un warning et sont ignorés.

    Args:
        folder_path: Chemin du dossier à parcourir.

    Returns:
        Liste consolidée de tous les Documents extraits.
    """
    folder = Path(folder_path)
    if not folder.is_dir():
        logger.error("Dossier introuvable : '%s'.", folder_path)
        return []

    all_docs: list[Document] = []
    candidates = [f for f in folder.iterdir() if f.is_file() and f.suffix.lower() in _PARSERS]

    if not candidates:
        logger.warning("Aucun fichier supporté trouvé dans '%s'.", folder_path)
        return []

    logger.info(
        "batch_parse — dossier='%s' | %d fichier(s) à traiter.",
        folder_path,
        len(candidates),
    )
    for file in candidates:
        docs = parse_file(str(file))
        all_docs.extend(docs)

    logger.info(
        "batch_parse terminé — %d fichier(s), %d Document(s) au total.",
        len(candidates),
        len(all_docs),
    )
    return all_docs


# ---------------------------------------------------------------------------
# Tests basiques — exécuter avec : python ingestion/file_parser.py
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    import os
    import sys
    import tempfile

    sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))
    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")

    def test_parse_txt() -> None:
        print("\n=== TEST parse_file(.txt) ===")
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", encoding="utf-8", delete=False) as f:
            f.write("Ceci est un document de test.\nLigne deux.")
            tmp = f.name
        try:
            docs = parse_file(tmp)
            assert len(docs) == 1, f"Attendu 1 Document, obtenu {len(docs)}"
            assert "document de test" in docs[0].page_content
            assert docs[0].metadata["file_type"] == "txt"
            assert docs[0].metadata["page"] == 1
            assert docs[0].metadata["filename"].endswith(".txt")
            print(f"  page_content : {docs[0].page_content[:60]}")
            print(f"  metadata     : {docs[0].metadata}")
            print("OK")
        finally:
            os.unlink(tmp)

    def test_unsupported_format() -> None:
        print("\n=== TEST format non supporté (.xml) ===")
        with tempfile.NamedTemporaryFile(suffix=".xml", delete=False) as f:
            f.write(b"<root><item>test</item></root>")
            tmp = f.name
        try:
            docs = parse_file(tmp)
            assert docs == [], f"Attendu [], obtenu {docs}"
            print("OK — retourne [] comme attendu")
        finally:
            os.unlink(tmp)

    def test_file_not_found() -> None:
        print("\n=== TEST fichier introuvable ===")
        docs = parse_file("/tmp/fichier_inexistant_telko.pdf")
        assert docs == [], f"Attendu [], obtenu {docs}"
        print("OK — retourne [] comme attendu")

    def test_empty_txt() -> None:
        print("\n=== TEST fichier TXT vide ===")
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False) as f:
            f.write("   \n  ")
            tmp = f.name
        try:
            docs = parse_file(tmp)
            assert docs == [], f"Attendu [], obtenu {docs}"
            print("OK — retourne [] comme attendu")
        finally:
            os.unlink(tmp)

    def test_batch_parse() -> None:
        print("\n=== TEST batch_parse() ===")
        with tempfile.TemporaryDirectory() as tmpdir:
            for i in range(3):
                Path(tmpdir, f"doc_{i}.txt").write_text(f"Contenu du document {i}.", encoding="utf-8")
            # Fichier ignoré (format non supporté)
            Path(tmpdir, "data.xml").write_bytes(b"<root/>")

            docs = batch_parse(tmpdir)
            assert len(docs) == 3, f"Attendu 3 Documents, obtenu {len(docs)}"
            print(f"  {len(docs)} Document(s) extraits (fichier .xml ignoré)")
            print("OK")

    def test_batch_empty_folder() -> None:
        print("\n=== TEST batch_parse() — dossier vide ===")
        with tempfile.TemporaryDirectory() as tmpdir:
            docs = batch_parse(tmpdir)
            assert docs == []
            print("OK — retourne [] comme attendu")

    def test_metadata_keys() -> None:
        print("\n=== TEST metadata complètes ===")
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", encoding="utf-8", delete=False) as f:
            f.write("Texte de vérification des metadata.")
            tmp = f.name
        try:
            docs = parse_file(tmp)
            meta = docs[0].metadata
            for key in ("source", "filename", "page", "file_type", "last_modified"):
                assert key in meta, f"Clé manquante dans metadata : '{key}'"
            assert meta["last_modified"] != "", "last_modified ne doit pas être vide"
            print(f"  Metadata : {meta}")
            print("OK")
        finally:
            os.unlink(tmp)

    # PDF / DOCX / PPTX nécessitent les libs et des fichiers réels :
    # Décommentez les lignes suivantes si vous avez des fichiers de test.
    # print("\n=== TEST parse_file(.pdf) ===")
    # docs = parse_file("/tmp/sample.pdf")
    # print(f"  {len(docs)} page(s) extraites")
    #
    # print("\n=== TEST parse_file(.docx) ===")
    # docs = parse_file("/tmp/sample.docx")
    # print(f"  {len(docs)} bloc(s) extraits")
    #
    # print("\n=== TEST parse_file(.pptx) ===")
    # docs = parse_file("/tmp/sample.pptx")
    # print(f"  {len(docs)} slide(s) extraites")

    test_parse_txt()
    test_unsupported_format()
    test_file_not_found()
    test_empty_txt()
    test_batch_parse()
    test_batch_empty_folder()
    test_metadata_keys()
    print("\nTous les tests sont passés.")
